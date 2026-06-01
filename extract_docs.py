# -*- coding: utf-8 -*-
"""
文件文字擷取工具（Phase A：結構化傾印）

讀取 files/ 內的 PDF / PPTX / DOCX，把文字連同結構標記
（DOCX 段落樣式、PPTX 投影片分頁、PDF 頁碼）輸出到 doc/_raw/，
供後續整理成正式 Markdown 使用。所有輸出皆為 UTF-8，確保繁體中文正常顯示。
"""
import sys
from pathlib import Path

# 確保在 Windows 終端機輸出繁體中文不會因 cp1252 編碼而崩潰
sys.stdout.reconfigure(encoding="utf-8")

import pdfplumber
from docx import Document
from docx.document import Document as _DocumentClass
from docx.oxml.ns import qn
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation

BASE = Path(__file__).resolve().parent
SRC = BASE / "files"
RAW = BASE / "doc" / "_raw"


def md_table(rows):
    """把二維字串陣列轉成 Markdown 表格。"""
    rows = [[(c or "").replace("\n", " ").strip() for c in r] for r in rows]
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = ["| " + " | ".join(rows[0]) + " |",
           "| " + " | ".join(["---"] * width) + " |"]
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


# ---------- DOCX ----------
def iter_block_items(parent):
    """依文件順序產生段落與表格（python-docx 預設不保證順序）。"""
    parent_elm = parent.element.body if isinstance(parent, _DocumentClass) else parent._tc
    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def extract_docx(path):
    doc = Document(str(path))
    lines = []
    for block in iter_block_items(doc):
        if isinstance(block, Table):
            rows = [[cell.text for cell in row.cells] for row in block.rows]
            lines.append("[TABLE]")
            lines.append(md_table(rows))
            lines.append("")
            continue
        text = block.text.strip()
        style = (block.style.name or "").strip() if block.style else ""
        if not text:
            continue
        # 偵測項目符號 / 編號清單
        numpr = block._p.find(qn("w:pPr"))
        is_list = numpr is not None and numpr.find(qn("w:numPr")) is not None
        tag = f"[{style}]" if style and style.lower() != "normal" else "[P]"
        if is_list and tag == "[P]":
            tag = "[LIST]"
        lines.append(f"{tag} {text}")
    return "\n".join(lines)


# ---------- PPTX ----------
def extract_pptx(path):
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"===== [SLIDE {i}] =====")
        title = slide.shapes.title
        title_text = title.text.strip() if title and title.text else ""
        if title_text:
            lines.append(f"[TITLE] {title_text}")
        for shape in slide.shapes:
            if shape == title:
                continue
            if shape.has_table:
                rows = [[c.text for c in row.cells] for row in shape.table.rows]
                lines.append("[TABLE]")
                lines.append(md_table(rows))
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs).strip() or para.text.strip()
                if not t:
                    continue
                lines.append(f"[BODY L{para.level}] {t}")
        lines.append("")
    return "\n".join(lines)


# ---------- PDF ----------
def extract_pdf(path):
    lines = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            lines.append(f"===== [PAGE {i}] =====")
            txt = page.extract_text(x_tolerance=1.5, y_tolerance=3) or ""
            lines.append(txt.strip())
            tables = page.extract_tables()
            for t in tables:
                lines.append("[TABLE]")
                lines.append(md_table(t))
            lines.append("")
    return "\n".join(lines)


def main():
    RAW.mkdir(parents=True, exist_ok=True)
    handlers = {".pdf": extract_pdf, ".pptx": extract_pptx, ".docx": extract_docx}
    for f in sorted(SRC.iterdir()):
        if f.suffix.lower() == ".md":
            print(f"SKIP (already markdown): {f.name}")
            continue
        handler = handlers.get(f.suffix.lower())
        if not handler:
            print(f"SKIP (unsupported): {f.name}")
            continue
        try:
            content = handler(f)
        except Exception as exc:  # noqa: BLE001
            print(f"ERROR {f.name}: {exc}")
            continue
        out = RAW / (f.stem + ".txt")
        out.write_text(content, encoding="utf-8")
        print(f"OK  {f.name} -> {out.relative_to(BASE)} ({len(content)} chars)")


if __name__ == "__main__":
    main()
