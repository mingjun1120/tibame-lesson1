"""
ShareBox 專案簡報 — 深色科技感 + 漸層版型
產出：files/專案簡報.pptx（覆蓋舊檔）

執行：python3 files/build_pptx.py
"""

from __future__ import annotations

from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ─────────────────────────────────────────────────────────────────────────────
# 設計系統
# ─────────────────────────────────────────────────────────────────────────────

# 色彩
BG_DEEP        = RGBColor(0x0A, 0x0E, 0x27)  # 主背景
BG_PANEL       = RGBColor(0x14, 0x19, 0x3D)  # 卡片底
BG_PANEL_HI    = RGBColor(0x1C, 0x22, 0x4D)  # 卡片高亮
LINE_SOFT      = RGBColor(0x2A, 0x32, 0x5C)  # 邊框

CYAN           = RGBColor(0x00, 0xD4, 0xFF)
PURPLE         = RGBColor(0x7C, 0x3A, 0xED)  # 深紫：用於卡片底、結構色
PURPLE_BRIGHT  = RGBColor(0xB1, 0x9C, 0xFF)  # 亮紫：用於漸層字（提升對比）
PINK           = RGBColor(0xFF, 0x6B, 0x9D)
LIME           = RGBColor(0xA3, 0xE6, 0x35)
AMBER          = RGBColor(0xFB, 0xBF, 0x24)

TEXT_PRIMARY   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SECONDARY = RGBColor(0xB8, 0xC5, 0xD6)
TEXT_MUTED     = RGBColor(0x7B, 0x86, 0xA8)
TEXT_DIM       = RGBColor(0x4F, 0x5A, 0x7A)

# 字型（macOS / Win 通吃；繁中走系統 fallback）
FONT_HEAD = "PingFang TC"
FONT_BODY = "PingFang TC"
FONT_MONO = "SF Mono"

# 投影片大小（16:9）
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# XML 小工具：python-pptx 沒有原生漸層 API，要直接寫 XML
# ─────────────────────────────────────────────────────────────────────────────

def _solid_fill_xml(rgb: RGBColor, alpha: int | None = None) -> str:
    """Solid fill, alpha 為百萬分比 (0-100000)；None 表示不透明"""
    a = f'<a:alpha val="{alpha}"/>' if alpha is not None else ""
    return (
        f'<a:solidFill><a:srgbClr val="{rgb}">{a}</a:srgbClr></a:solidFill>'
    )


def _apply_gradient_fill(shape, stops: list[tuple[int, RGBColor, int | None]], angle_deg: int = 45):
    """
    Apply linear gradient. stops = list of (position_percent, color, alpha_or_None).
    position_percent: 0~100 -> 0~100000
    angle_deg: gradient direction (PowerPoint angle, 60000 units per degree)
    """
    sp = shape.fill._xPr  # ShapeProperties
    # 清掉現有的 fill
    for tag in ("a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:noFill"):
        for el in sp.findall(qn(tag)):
            sp.remove(el)

    stop_xml = ""
    for pos, color, alpha in stops:
        a = f'<a:alpha val="{alpha}"/>' if alpha is not None else ""
        stop_xml += (
            f'<a:gs pos="{pos * 1000}">'
            f'<a:srgbClr val="{color}">{a}</a:srgbClr>'
            f'</a:gs>'
        )

    grad_xml = (
        f'<a:gradFill flip="none" rotWithShape="1" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:gsLst>{stop_xml}</a:gsLst>'
        f'<a:lin ang="{angle_deg * 60000}" scaled="0"/>'
        f'<a:tileRect/>'
        f'</a:gradFill>'
    )
    sp.insert(
        list(sp).index(sp.find(qn("a:xfrm"))) + 1 if sp.find(qn("a:xfrm")) is not None else 0,
        etree.fromstring(grad_xml),
    )


def _no_line(shape):
    ln = shape.line
    ln.fill.background()


def _set_line(shape, rgb: RGBColor, width_pt: float = 1.0, alpha: int | None = None):
    ln = shape.line
    ln.color.rgb = rgb
    ln.width = Pt(width_pt)
    if alpha is not None:
        # 加 alpha 到 line color
        ln_xml = shape.line._get_or_add_ln()
        solid = ln_xml.find(qn("a:solidFill"))
        if solid is not None:
            ln_xml.remove(solid)
        ln_xml.insert(0, etree.fromstring(
            f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            f'<a:srgbClr val="{rgb}"><a:alpha val="{alpha}"/></a:srgbClr>'
            f'</a:solidFill>'
        ))


def _set_shadow(shape, blur_pt: float = 20, distance_pt: float = 6, alpha: int = 40000):
    """在 shape 加柔和陰影，讓深色卡片有浮起感"""
    sp = shape.fill._xPr
    # 移除舊的 effectLst
    for tag in ("a:effectLst",):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    shadow = (
        f'<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{int(blur_pt * 12700)}" dist="{int(distance_pt * 12700)}" '
        f'dir="5400000" algn="t" rotWithShape="0">'
        f'<a:srgbClr val="000000"><a:alpha val="{alpha}"/></a:srgbClr>'
        f'</a:outerShdw>'
        f'</a:effectLst>'
    )
    sp.append(etree.fromstring(shadow))


# ─────────────────────────────────────────────────────────────────────────────
# 文字小工具
# ─────────────────────────────────────────────────────────────────────────────

def add_text(
    slide, left, top, width, height,
    text: str,
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = TEXT_PRIMARY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font: str = FONT_BODY,
    letter_spacing: int | None = None,  # in 1/100 pt, e.g. 200 = 2pt
    line_spacing: float | None = None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        # 中文字型 fallback
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)
        cs = rPr.find(qn("a:cs"))
        if cs is None:
            cs = etree.SubElement(rPr, qn("a:cs"))
        cs.set("typeface", font)

        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

        if letter_spacing is not None:
            rPr.set("spc", str(letter_spacing))
    return tb


def add_gradient_text(
    slide, left, top, width, height,
    text: str,
    *,
    size: int = 80,
    bold: bool = True,
    stops: list[tuple[int, RGBColor]] = None,
    angle_deg: int = 0,
    align=PP_ALIGN.LEFT,
    font: str = FONT_HEAD,
    line_spacing: float | None = None,
):
    """把文字填上漸層"""
    if stops is None:
        stops = [(0, CYAN), (100, PURPLE)]
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold

    rPr = run._r.get_or_add_rPr()
    # 中文 fallback
    for ns in ("a:ea", "a:cs"):
        el = rPr.find(qn(ns))
        if el is None:
            el = etree.SubElement(rPr, qn(ns))
        el.set("typeface", font)

    # 清掉預設 solidFill
    for tag in ("a:solidFill", "a:gradFill", "a:noFill"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)

    stop_xml = "".join(
        f'<a:gs pos="{pos * 1000}"><a:srgbClr val="{color}"/></a:gs>'
        for pos, color in stops
    )
    grad = (
        f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:gsLst>{stop_xml}</a:gsLst>'
        f'<a:lin ang="{angle_deg * 60000}" scaled="0"/>'
        f'</a:gradFill>'
    )
    # gradFill 必須在 rPr 子節點中正確順序：在 ln 之後、effectLst 之前
    # 簡單起見直接 append；PowerPoint 對 run 屬性順序較寬容
    rPr.append(etree.fromstring(grad))
    return tb


# ─────────────────────────────────────────────────────────────────────────────
# 共用元件
# ─────────────────────────────────────────────────────────────────────────────

def paint_background(slide):
    """畫面層級背景：深色 + 左上漸層光暈 + 右下漸層光暈"""
    # 全幅底色矩形
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DEEP
    _no_line(bg)

    # 左上紫色光暈
    glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.5), Inches(-2.5), Inches(6), Inches(6))
    _apply_gradient_fill(
        glow1,
        stops=[(0, PURPLE, 35000), (100, PURPLE, 0)],
        angle_deg=135,
    )
    _no_line(glow1)

    # 右下青色光暈
    glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(4.5), Inches(6.5), Inches(6.5))
    _apply_gradient_fill(
        glow2,
        stops=[(0, CYAN, 30000), (100, CYAN, 0)],
        angle_deg=315,
    )
    _no_line(glow2)

    # 微噪點：用一條細漸層線當 grid 提示
    grid_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, Inches(7.42), SLIDE_W, Inches(0.02),
    )
    _apply_gradient_fill(
        grid_line,
        stops=[(0, CYAN, 0), (50, CYAN, 60000), (100, PURPLE, 0)],
        angle_deg=0,
    )
    _no_line(grid_line)


def add_eyebrow(slide, left, top, text: str, color: RGBColor = CYAN):
    """小標 chip：CYAN 字 + 字距拉開"""
    add_text(
        slide, left, top, Inches(5), Inches(0.3),
        text.upper(), size=12, bold=True, color=color,
        font="Inter", letter_spacing=400,
    )


def add_page_number(slide, idx: int, total: int):
    add_text(
        slide,
        Inches(12.4), Inches(6.95), Inches(1), Inches(0.3),
        f"{idx:02d} / {total:02d}", size=10, color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT, font="Inter",
    )
    add_text(
        slide,
        Inches(0.5), Inches(6.95), Inches(4), Inches(0.3),
        "SHAREBOX · 專案簡報", size=10, color=TEXT_MUTED,
        align=PP_ALIGN.LEFT, font="Inter", letter_spacing=300,
    )


def add_card(
    slide, left, top, width, height,
    *,
    accent: RGBColor = CYAN,
    radius_pt: float = 12.0,
):
    """加一張深色卡片（左側帶 accent 條）"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # 改圓角為固定大小
    card.adjustments[0] = 0.08
    card.fill.solid()
    card.fill.fore_color.rgb = BG_PANEL
    _set_line(card, LINE_SOFT, 0.75)
    _set_shadow(card, blur_pt=24, distance_pt=8, alpha=50000)

    # accent 細條
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + Inches(0.18), top + Inches(0.25),
        Inches(0.06), Inches(0.5),
    )
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    _no_line(bar)
    return card


def add_number_badge(slide, left, top, n: int, color: RGBColor = CYAN):
    """大編號數字（透明感）"""
    add_text(
        slide, left, top, Inches(1.5), Inches(1.2),
        f"{n:02d}", size=48, bold=True, color=color, font="Inter",
        letter_spacing=-50,
    )


# ─────────────────────────────────────────────────────────────────────────────
# 各頁
# ─────────────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)

    # 大型背景幾何
    blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(-1), Inches(9), Inches(9))
    _apply_gradient_fill(
        blob,
        stops=[(0, PURPLE, 25000), (60, CYAN, 12000), (100, BG_DEEP, 0)],
        angle_deg=45,
    )
    _no_line(blob)

    # 右下小圓
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(5.5), Inches(3), Inches(3))
    _apply_gradient_fill(
        blob2,
        stops=[(0, PINK, 35000), (100, PINK, 0)],
        angle_deg=90,
    )
    _no_line(blob2)

    # 品牌標記
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(0.6), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CYAN
    _no_line(dot)
    add_text(
        slide, Inches(0.9), Inches(0.55), Inches(3), Inches(0.3),
        "SHAREBOX", size=13, bold=True, color=TEXT_PRIMARY,
        font="Inter", letter_spacing=400,
    )

    # 小副標
    add_text(
        slide, Inches(0.6), Inches(2.0), Inches(4), Inches(0.4),
        "COMMUNITY · SHARING · TRUST", size=11, bold=True, color=CYAN,
        font="Inter", letter_spacing=600,
    )

    # 主標
    add_gradient_text(
        slide, Inches(0.6), Inches(2.4), Inches(10), Inches(2.0),
        "ShareBox",
        size=110, bold=True,
        stops=[(0, CYAN), (60, PURPLE_BRIGHT), (100, PINK)],
        angle_deg=15,
        line_spacing=1.0,
    )

    # 副標題
    add_text(
        slide, Inches(0.65), Inches(4.4), Inches(10), Inches(0.8),
        "社區共享工具平台", size=36, bold=True, color=TEXT_PRIMARY,
    )

    # tagline
    add_text(
        slide, Inches(0.65), Inches(5.25), Inches(10), Inches(0.5),
        "讓鄰居的閒置工具動起來", size=20, color=TEXT_SECONDARY,
    )

    # 漸層分隔線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(6.15), Inches(2.5), Inches(0.04))
    _apply_gradient_fill(line, stops=[(0, CYAN, 100000), (100, PURPLE, 0)], angle_deg=0)
    _no_line(line)

    # 右下日期
    add_text(
        slide, Inches(0.65), Inches(6.5), Inches(8), Inches(0.4),
        "專案簡報 · 2026 年 1 月", size=14, color=TEXT_MUTED,
        font="Inter", letter_spacing=200,
    )


def slide_problem(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "01 · The Problem", CYAN)
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(12), Inches(1.0),
        "我們要解決的問題", size=44, bold=True, color=TEXT_PRIMARY,
    )

    items = [
        ("01", "🔧", "閒置率高", "家用工具買了卻很少用 ——\n一年常使用不到 3 次", CYAN),
        ("02", "📦", "空間浪費", "工具佔據居家收納空間\n造成資源閒置與浪費", PURPLE),
        ("03", "🤝", "缺乏信任", "鄰居想互借卻缺乏信任\n與損壞賠償的保障", PINK),
        ("04", "🏪", "租借不便", "租借店租金高、據點少、\n品項受限，且需親自往返", AMBER),
    ]

    card_w = Inches(5.7)
    card_h = Inches(2.35)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    start_x = Inches(0.6)
    start_y = Inches(2.3)

    for i, (num, icon, title, body, color) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_card(slide, x, y, card_w, card_h, accent=color)
        # 大圖示
        add_text(
            slide, x + Inches(0.45), y + Inches(0.3), Inches(1), Inches(0.8),
            icon, size=36, color=color,
        )
        # 編號
        add_text(
            slide, x + card_w - Inches(1.0), y + Inches(0.3), Inches(0.8), Inches(0.5),
            num, size=14, bold=True, color=TEXT_DIM, font="Inter",
            align=PP_ALIGN.RIGHT, letter_spacing=300,
        )
        # 標題
        add_text(
            slide, x + Inches(0.45), y + Inches(1.05), card_w - Inches(0.9), Inches(0.5),
            title, size=22, bold=True, color=TEXT_PRIMARY,
        )
        # 描述
        add_text(
            slide, x + Inches(0.45), y + Inches(1.55), card_w - Inches(0.9), Inches(0.7),
            body, size=13, color=TEXT_SECONDARY, line_spacing=1.4,
        )


def slide_solution(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "02 · Our Solution", PINK)

    add_gradient_text(
        slide, Inches(0.6), Inches(0.9), Inches(8), Inches(1.4),
        "ShareBox",
        size=70, bold=True,
        stops=[(0, CYAN), (100, PURPLE_BRIGHT)],
        angle_deg=15,
    )
    add_text(
        slide, Inches(0.6), Inches(2.05), Inches(8), Inches(0.7),
        "以「社區」為單位的工具共享 App", size=26, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, Inches(0.6), Inches(2.7), Inches(7.2), Inches(1.2),
        "把鄰居的閒置工具，轉化為社區共用資源；\n用信任機制與押金代管，補上鄰里互借的最後一塊拼圖。",
        size=14, color=TEXT_SECONDARY, line_spacing=1.6,
    )

    # 4 個 highlight bullets
    bullets = [
        ("地圖搜尋", "附近鄰居的閒置工具一目了然", CYAN),
        ("即時聊天", "App 內媒合、約定借用與歸還", PURPLE),
        ("押金代管", "雙向評價，建立鄰里信任", PINK),
        ("綠界金流", "押金第三方代管，安全有保障", LIME),
    ]
    for i, (title, body, color) in enumerate(bullets):
        y = Inches(4.2 + i * 0.65)
        # 圓點
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.15), Inches(0.18), Inches(0.18))
        d.fill.solid()
        d.fill.fore_color.rgb = color
        _no_line(d)
        # 標題
        add_text(
            slide, Inches(0.9), y, Inches(2.5), Inches(0.45),
            title, size=16, bold=True, color=TEXT_PRIMARY,
        )
        # 描述
        add_text(
            slide, Inches(3.0), y, Inches(5.5), Inches(0.45),
            body, size=14, color=TEXT_SECONDARY,
        )

    # 右側：手機概念圖（用幾何形狀疊出）
    phone_x = Inches(9.2)
    phone_y = Inches(1.3)
    phone_w = Inches(3.1)
    phone_h = Inches(5.5)

    # 手機外殼
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, phone_x, phone_y, phone_w, phone_h)
    phone.adjustments[0] = 0.10
    _apply_gradient_fill(
        phone,
        stops=[(0, BG_PANEL_HI, 100000), (100, BG_DEEP, 100000)],
        angle_deg=135,
    )
    _set_line(phone, CYAN, 1.2, alpha=60000)
    _set_shadow(phone, blur_pt=30, distance_pt=12, alpha=60000)

    # 螢幕內容：頂部 nav
    nav = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.2), phone_y + Inches(0.35),
        phone_w - Inches(0.4), Inches(0.5),
    )
    nav.adjustments[0] = 0.3
    nav.fill.solid()
    nav.fill.fore_color.rgb = BG_DEEP
    _no_line(nav)
    add_text(
        slide, phone_x + Inches(0.4), phone_y + Inches(0.43), Inches(2.5), Inches(0.3),
        "附近工具", size=12, bold=True, color=TEXT_PRIMARY,
    )
    # 漸層 chip
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + phone_w - Inches(0.85), phone_y + Inches(0.43),
        Inches(0.5), Inches(0.32),
    )
    chip.adjustments[0] = 0.5
    _apply_gradient_fill(chip, stops=[(0, CYAN, 100000), (100, PURPLE, 100000)], angle_deg=0)
    _no_line(chip)

    # 卡片 1（高亮）
    card1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.2), phone_y + Inches(1.05),
        phone_w - Inches(0.4), Inches(1.2),
    )
    card1.adjustments[0] = 0.15
    _apply_gradient_fill(card1, stops=[(0, CYAN, 25000), (100, PURPLE, 25000)], angle_deg=135)
    _set_line(card1, CYAN, 0.5, alpha=80000)
    # icon dot
    ic1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.4), phone_y + Inches(1.25),
        Inches(0.6), Inches(0.6),
    )
    ic1.adjustments[0] = 0.3
    ic1.fill.solid()
    ic1.fill.fore_color.rgb = BG_DEEP
    _no_line(ic1)
    add_text(
        slide, phone_x + Inches(0.4), phone_y + Inches(1.28), Inches(0.6), Inches(0.5),
        "🔨", size=22, color=CYAN, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, phone_x + Inches(1.1), phone_y + Inches(1.25), Inches(1.8), Inches(0.3),
        "電鑽 · 5F", size=11, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, phone_x + Inches(1.1), phone_y + Inches(1.55), Inches(1.8), Inches(0.3),
        "距離 12m", size=9, color=TEXT_MUTED,
    )
    # 押金金額
    add_text(
        slide, phone_x + Inches(1.1), phone_y + Inches(1.82), Inches(1.8), Inches(0.3),
        "押金 300 元", size=10, bold=True, color=CYAN,
    )

    # 卡片 2
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.2), phone_y + Inches(2.4),
        phone_w - Inches(0.4), Inches(0.9),
    )
    card2.adjustments[0] = 0.15
    card2.fill.solid()
    card2.fill.fore_color.rgb = BG_PANEL_HI
    _set_line(card2, LINE_SOFT, 0.5)
    ic2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.4), phone_y + Inches(2.55),
        Inches(0.5), Inches(0.5),
    )
    ic2.adjustments[0] = 0.3
    ic2.fill.solid()
    ic2.fill.fore_color.rgb = BG_DEEP
    _no_line(ic2)
    add_text(
        slide, phone_x + Inches(0.4), phone_y + Inches(2.55), Inches(0.5), Inches(0.4),
        "🪜", size=18, color=PURPLE, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, phone_x + Inches(1.0), phone_y + Inches(2.55), Inches(1.8), Inches(0.3),
        "鋁梯 · 3F", size=11, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, phone_x + Inches(1.0), phone_y + Inches(2.85), Inches(1.8), Inches(0.3),
        "距離 35m · 押金 200", size=9, color=TEXT_MUTED,
    )

    # 卡片 3
    card3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.2), phone_y + Inches(3.45),
        phone_w - Inches(0.4), Inches(0.9),
    )
    card3.adjustments[0] = 0.15
    card3.fill.solid()
    card3.fill.fore_color.rgb = BG_PANEL_HI
    _set_line(card3, LINE_SOFT, 0.5)
    ic3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.4), phone_y + Inches(3.6),
        Inches(0.5), Inches(0.5),
    )
    ic3.adjustments[0] = 0.3
    ic3.fill.solid()
    ic3.fill.fore_color.rgb = BG_DEEP
    _no_line(ic3)
    add_text(
        slide, phone_x + Inches(0.4), phone_y + Inches(3.6), Inches(0.5), Inches(0.4),
        "🧰", size=18, color=PINK, align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, phone_x + Inches(1.0), phone_y + Inches(3.6), Inches(1.8), Inches(0.3),
        "扳手組 · 8F", size=11, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, phone_x + Inches(1.0), phone_y + Inches(3.9), Inches(1.8), Inches(0.3),
        "距離 48m · 押金 150", size=9, color=TEXT_MUTED,
    )

    # 底部按鈕
    btn = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        phone_x + Inches(0.4), phone_y + Inches(4.65),
        phone_w - Inches(0.8), Inches(0.55),
    )
    btn.adjustments[0] = 0.5
    _apply_gradient_fill(btn, stops=[(0, CYAN, 100000), (100, PURPLE, 100000)], angle_deg=0)
    _no_line(btn)
    add_text(
        slide, phone_x + Inches(0.4), phone_y + Inches(4.78), phone_w - Inches(0.8), Inches(0.4),
        "立即申請借用 →", size=13, bold=True, color=TEXT_PRIMARY,
        align=PP_ALIGN.CENTER,
    )


def slide_features(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "03 · Core Features", PURPLE)
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(10), Inches(1.0),
        "核心功能", size=44, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, Inches(0.6), Inches(1.65), Inches(10), Inches(0.4),
        "六大模組，覆蓋從刊登、媒合、金流到歸還的完整流程",
        size=15, color=TEXT_SECONDARY,
    )

    features = [
        ("📸", "工具刊登", "拍照、選分類、設定押金", CYAN),
        ("🗺️", "地圖與清單", "附近工具一鍵搜尋", PURPLE),
        ("💬", "即時聊天", "借用申請與訂單溝通", PINK),
        ("💳", "押金代管", "綠界 ECPay 第三方金流", LIME),
        ("✅", "歸還與評價", "雙向評價建立信任", AMBER),
        ("🔔", "推播通知", "申請、同意、歸還提醒", CYAN),
    ]

    card_w = Inches(4.0)
    card_h = Inches(2.2)
    gap = Inches(0.18)
    start_x = Inches(0.6)
    start_y = Inches(2.55)

    for i, (icon, title, body, color) in enumerate(features):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        add_card(slide, x, y, card_w, card_h, accent=color)

        # icon 容器
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.45), y + Inches(0.35),
            Inches(0.85), Inches(0.85),
        )
        icon_box.adjustments[0] = 0.25
        _apply_gradient_fill(
            icon_box,
            stops=[(0, color, 35000), (100, color, 10000)],
            angle_deg=135,
        )
        _set_line(icon_box, color, 0.5, alpha=50000)
        add_text(
            slide, x + Inches(0.45), y + Inches(0.38), Inches(0.85), Inches(0.8),
            icon, size=30, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER,
        )

        # 編號
        add_text(
            slide, x + card_w - Inches(0.9), y + Inches(0.4), Inches(0.6), Inches(0.4),
            f"0{i+1}", size=14, bold=True, color=TEXT_DIM, font="Inter",
            align=PP_ALIGN.RIGHT, letter_spacing=300,
        )

        # 標題
        add_text(
            slide, x + Inches(0.45), y + Inches(1.3), card_w - Inches(0.9), Inches(0.45),
            title, size=20, bold=True, color=TEXT_PRIMARY,
        )
        # 描述
        add_text(
            slide, x + Inches(0.45), y + Inches(1.72), card_w - Inches(0.9), Inches(0.4),
            body, size=12, color=TEXT_SECONDARY,
        )


def slide_market(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "04 · Market & Business", LIME)
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(10), Inches(1.0),
        "目標市場與商業模式", size=42, bold=True, color=TEXT_PRIMARY,
    )

    # 左半：巨大數字
    add_text(
        slide, Inches(0.6), Inches(2.5), Inches(6), Inches(0.5),
        "TOTAL ADDRESSABLE MARKET", size=12, bold=True, color=CYAN,
        font="Inter", letter_spacing=400,
    )
    add_gradient_text(
        slide, Inches(0.6), Inches(2.9), Inches(7), Inches(2.5),
        "3.5M",
        size=180, bold=True,
        stops=[(0, CYAN), (100, PURPLE_BRIGHT)],
        angle_deg=15,
        line_spacing=0.9,
    )
    add_text(
        slide, Inches(0.6), Inches(5.5), Inches(6), Inches(0.5),
        "都會區公寓大廈住戶", size=22, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, Inches(0.6), Inches(6.0), Inches(6), Inches(0.4),
        "市場尚無直接競品，以信任機制為核心差異化",
        size=13, color=TEXT_SECONDARY,
    )

    # 右半：商業模式
    add_text(
        slide, Inches(8.0), Inches(2.5), Inches(5), Inches(0.4),
        "REVENUE STREAMS", size=12, bold=True, color=PINK,
        font="Inter", letter_spacing=400,
    )

    items = [
        ("免費版", "降低使用門檻", "免費", CYAN),
        ("Premium 訂閱", "進階篩選與優先媒合", "NT$ 99/月", PURPLE),
        ("交易抽成", "每筆付費借用", "10%", PINK),
    ]
    for i, (name, desc, price, color) in enumerate(items):
        y = Inches(3.0 + i * 1.25)
        add_card(slide, Inches(8.0), y, Inches(4.8), Inches(1.05), accent=color)
        add_text(
            slide, Inches(8.4), y + Inches(0.18), Inches(2.5), Inches(0.4),
            name, size=17, bold=True, color=TEXT_PRIMARY,
        )
        add_text(
            slide, Inches(8.4), y + Inches(0.6), Inches(3), Inches(0.35),
            desc, size=12, color=TEXT_SECONDARY,
        )
        # 價格
        add_text(
            slide, Inches(11.2), y + Inches(0.32), Inches(1.5), Inches(0.5),
            price, size=18, bold=True, color=color, font="Inter",
            align=PP_ALIGN.RIGHT,
        )


def slide_team(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "05 · The Team", PINK)
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(10), Inches(1.0),
        "專案團隊", size=44, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, Inches(0.6), Inches(1.65), Inches(10), Inches(0.4),
        "5 人核心團隊，橫跨產品、技術、設計與行銷",
        size=15, color=TEXT_SECONDARY,
    )

    team = [
        ("林佳穎", "Lin Chia-Ying", "專案經理", "Project Manager", "佳", CYAN),
        ("陳冠宇", "Chen Kuan-Yu", "後端工程師", "Backend Engineer", "宇", PURPLE),
        ("張哲瑋", "Chang Che-Wei", "前端 / App", "Frontend / App", "瑋", PINK),
        ("吳采潔", "Wu Tsai-Chieh", "UI/UX 設計師", "Product Designer", "潔", LIME),
        ("黃柏翔", "Huang Po-Hsiang", "行銷企劃", "Marketing", "翔", AMBER),
    ]

    card_w = Inches(2.35)
    card_h = Inches(4.3)
    gap = Inches(0.13)
    start_x = Inches(0.6)
    start_y = Inches(2.5)

    for i, (name, name_en, role, role_en, char, color) in enumerate(team):
        x = start_x + i * (card_w + gap)
        y = start_y

        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.adjustments[0] = 0.08
        card.fill.solid()
        card.fill.fore_color.rgb = BG_PANEL
        _set_line(card, LINE_SOFT, 0.75)
        _set_shadow(card, blur_pt=24, distance_pt=8, alpha=50000)

        # 頭像背景（漸層圓）
        ava_size = Inches(1.6)
        ava_x = x + (card_w - ava_size) / 2
        ava_y = y + Inches(0.5)
        ava = slide.shapes.add_shape(MSO_SHAPE.OVAL, ava_x, ava_y, ava_size, ava_size)
        _apply_gradient_fill(
            ava,
            stops=[(0, color, 100000), (100, PURPLE, 80000)],
            angle_deg=135,
        )
        _no_line(ava)
        # 頭像中字
        add_text(
            slide, ava_x, ava_y + Inches(0.32), ava_size, Inches(1),
            char, size=46, bold=True, color=TEXT_PRIMARY,
            align=PP_ALIGN.CENTER,
        )

        # 中文姓名
        add_text(
            slide, x, y + Inches(2.3), card_w, Inches(0.5),
            name, size=22, bold=True, color=TEXT_PRIMARY,
            align=PP_ALIGN.CENTER,
        )
        # 英文姓名
        add_text(
            slide, x, y + Inches(2.78), card_w, Inches(0.3),
            name_en, size=10, color=TEXT_MUTED, font="Inter",
            align=PP_ALIGN.CENTER, letter_spacing=200,
        )

        # 分隔
        div = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + card_w / 2 - Inches(0.3), y + Inches(3.2),
            Inches(0.6), Inches(0.02),
        )
        div.fill.solid()
        div.fill.fore_color.rgb = color
        _no_line(div)

        # 職稱
        add_text(
            slide, x, y + Inches(3.4), card_w, Inches(0.4),
            role, size=14, bold=True, color=color,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide, x, y + Inches(3.78), card_w, Inches(0.3),
            role_en, size=10, color=TEXT_MUTED, font="Inter",
            align=PP_ALIGN.CENTER, letter_spacing=150,
        )


def slide_timeline(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "06 · Roadmap", CYAN)
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(10), Inches(1.0),
        "專案時程與里程碑", size=44, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, Inches(0.6), Inches(1.65), Inches(10), Inches(0.4),
        "從需求確認到正式上線，預計 8 個月完成",
        size=15, color=TEXT_SECONDARY,
    )

    # 時間軸主線（漸層）
    line_y = Inches(4.3)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), line_y, Inches(10.3), Inches(0.05),
    )
    _apply_gradient_fill(
        line,
        stops=[(0, CYAN, 100000), (33, PURPLE, 100000), (66, PINK, 100000), (100, AMBER, 100000)],
        angle_deg=0,
    )
    _no_line(line)

    milestones = [
        ("M1", "需求確認與設計", "2026", "1 – 2 月", "釐清功能、UI/UX 定稿", CYAN),
        ("M2", "MVP 開發", "2026", "3 – 5 月", "核心功能上線可內部測試", PURPLE),
        ("M3", "封閉測試 Beta", "2026", "6 – 7 月", "500 位社區住戶試用", PINK),
        ("M4", "正式上線", "2026", "8 月", "公開上架雙平台", AMBER),
    ]

    # 4 個里程碑點
    n = len(milestones)
    span_left = Inches(1.5)
    span_right = Inches(11.8)
    step = (span_right - span_left) / (n - 1)

    for i, (m, name, year, period, desc, color) in enumerate(milestones):
        cx = span_left + i * step
        # 外圈光暈
        glow = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            cx - Inches(0.4), line_y - Inches(0.375),
            Inches(0.8), Inches(0.8),
        )
        _apply_gradient_fill(
            glow,
            stops=[(0, color, 60000), (100, color, 0)],
            angle_deg=0,
        )
        _no_line(glow)
        # 內圓
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            cx - Inches(0.18), line_y - Inches(0.155),
            Inches(0.36), Inches(0.36),
        )
        _apply_gradient_fill(
            dot,
            stops=[(0, color, 100000), (100, PURPLE, 80000)],
            angle_deg=135,
        )
        _set_line(dot, TEXT_PRIMARY, 1.0, alpha=40000)

        # 上方：M 標籤 + 時間
        top_y = line_y - Inches(1.8)
        add_text(
            slide, cx - Inches(1.3), top_y, Inches(2.6), Inches(0.35),
            m, size=14, bold=True, color=color, font="Inter",
            align=PP_ALIGN.CENTER, letter_spacing=300,
        )
        add_text(
            slide, cx - Inches(1.3), top_y + Inches(0.4), Inches(2.6), Inches(0.45),
            name, size=18, bold=True, color=TEXT_PRIMARY,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide, cx - Inches(1.3), top_y + Inches(0.95), Inches(2.6), Inches(0.3),
            f"{year} · {period}", size=11, color=TEXT_MUTED, font="Inter",
            align=PP_ALIGN.CENTER, letter_spacing=200,
        )

        # 下方：描述
        bottom_y = line_y + Inches(0.65)
        add_text(
            slide, cx - Inches(1.4), bottom_y, Inches(2.8), Inches(0.6),
            desc, size=12, color=TEXT_SECONDARY,
            align=PP_ALIGN.CENTER, line_spacing=1.4,
        )


def slide_budget(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "07 · Budget", AMBER)
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(10), Inches(1.0),
        "專案預算概覽", size=44, bold=True, color=TEXT_PRIMARY,
    )

    # 左側：總預算
    add_text(
        slide, Inches(0.6), Inches(2.5), Inches(5), Inches(0.4),
        "TOTAL BUDGET", size=12, bold=True, color=AMBER,
        font="Inter", letter_spacing=500,
    )
    add_gradient_text(
        slide, Inches(0.6), Inches(2.95), Inches(6.4), Inches(2.0),
        "NT$ 4.8M",
        size=78, bold=True,
        stops=[(0, AMBER), (100, PINK)],
        angle_deg=15,
        line_spacing=1.0,
    )
    add_text(
        slide, Inches(0.6), Inches(5.0), Inches(6), Inches(0.5),
        "新台幣 480 萬元", size=24, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, Inches(0.6), Inches(5.6), Inches(6), Inches(0.4),
        "涵蓋人事、雲端、行銷與雜支",
        size=14, color=TEXT_SECONDARY,
    )

    # 右側：分項
    items = [
        ("人事費",       320, 66.7, CYAN),
        ("行銷推廣",      80, 16.7, PURPLE),
        ("雲端與基礎設施", 40,  8.3, PINK),
        ("其他與雜支",    40,  8.3, LIME),
    ]
    base_x = Inches(7.2)
    base_y = Inches(2.5)
    row_h = Inches(1.05)
    bar_max_w = Inches(5.4)

    for i, (name, amount, pct, color) in enumerate(items):
        y = base_y + i * row_h
        # 名稱
        add_text(
            slide, base_x, y, Inches(2.5), Inches(0.35),
            name, size=15, bold=True, color=TEXT_PRIMARY,
        )
        # 金額
        add_text(
            slide, base_x + Inches(2.5), y, Inches(1.6), Inches(0.35),
            f"NT$ {amount} 萬", size=13, bold=True, color=color, font="Inter",
        )
        # 百分比
        add_text(
            slide, Inches(11.7), y, Inches(1.1), Inches(0.35),
            f"{pct:.1f}%", size=13, bold=True, color=TEXT_MUTED, font="Inter",
            align=PP_ALIGN.RIGHT,
        )
        # bar 底
        bar_y = y + Inches(0.5)
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            base_x, bar_y, bar_max_w, Inches(0.18),
        )
        bar_bg.adjustments[0] = 0.5
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = BG_PANEL
        _no_line(bar_bg)
        # bar 前景
        fg_w = int(bar_max_w * (pct / 100.0))
        bar_fg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            base_x, bar_y, fg_w, Inches(0.18),
        )
        bar_fg.adjustments[0] = 0.5
        _apply_gradient_fill(bar_fg, stops=[(0, color, 100000), (100, PURPLE, 100000)], angle_deg=0)
        _no_line(bar_fg)


def slide_kpi(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_page_number(slide, idx, total)

    add_eyebrow(slide, Inches(0.6), Inches(0.55), "08 · Year 1 KPI", LIME)
    add_text(
        slide, Inches(0.6), Inches(0.85), Inches(10), Inches(1.0),
        "首年目標 KPI", size=44, bold=True, color=TEXT_PRIMARY,
    )
    add_text(
        slide, Inches(0.6), Inches(1.65), Inches(10), Inches(0.4),
        "上線一年內必須達成的關鍵指標",
        size=15, color=TEXT_SECONDARY,
    )

    kpis = [
        ("10K",  "累積下載數",       "Total Downloads",      CYAN),
        ("3K",   "月活躍用戶 MAU",   "Monthly Active Users", PURPLE),
        ("500",  "每月成功借用",      "Successful Borrowings",PINK),
        ("8%",   "Premium 轉換率",  "Paid Conversion Rate",  LIME),
    ]

    card_w = Inches(3.0)
    card_h = Inches(4.3)
    gap = Inches(0.12)
    start_x = Inches(0.6)
    start_y = Inches(2.55)

    for i, (num, label, label_en, color) in enumerate(kpis):
        x = start_x + i * (card_w + gap)
        y = start_y

        # 卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.adjustments[0] = 0.08
        _apply_gradient_fill(
            card,
            stops=[(0, BG_PANEL_HI, 100000), (100, BG_PANEL, 100000)],
            angle_deg=135,
        )
        _set_line(card, LINE_SOFT, 0.75)
        _set_shadow(card, blur_pt=24, distance_pt=8, alpha=50000)

        # 上方 color 條
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.5), y + Inches(0.35),
            Inches(0.6), Inches(0.06),
        )
        _apply_gradient_fill(bar, stops=[(0, color, 100000), (100, PURPLE, 100000)], angle_deg=0)
        _no_line(bar)

        # 巨大數字（漸層）
        add_gradient_text(
            slide, x + Inches(0.3), y + Inches(0.9), card_w - Inches(0.6), Inches(2.0),
            num, size=88, bold=True,
            stops=[(0, color), (100, PURPLE_BRIGHT)],
            angle_deg=15,
            align=PP_ALIGN.LEFT,
            line_spacing=0.95,
        )

        # 中文標籤
        add_text(
            slide, x + Inches(0.5), y + Inches(3.1), card_w - Inches(1), Inches(0.5),
            label, size=18, bold=True, color=TEXT_PRIMARY,
        )
        # 英文標籤
        add_text(
            slide, x + Inches(0.5), y + Inches(3.55), card_w - Inches(1), Inches(0.4),
            label_en, size=11, color=TEXT_MUTED, font="Inter",
            letter_spacing=200,
        )

        # KPI 編號
        add_text(
            slide, x + card_w - Inches(0.9), y + Inches(0.35), Inches(0.6), Inches(0.4),
            f"K{i+1}", size=12, bold=True, color=TEXT_DIM, font="Inter",
            align=PP_ALIGN.RIGHT, letter_spacing=300,
        )


def slide_closing(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)

    # 大型背景幾何
    blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(3), Inches(8), Inches(8))
    _apply_gradient_fill(
        blob,
        stops=[(0, CYAN, 30000), (100, CYAN, 0)],
        angle_deg=45,
    )
    _no_line(blob)
    blob2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-3), Inches(8), Inches(8))
    _apply_gradient_fill(
        blob2,
        stops=[(0, PINK, 30000), (100, PINK, 0)],
        angle_deg=135,
    )
    _no_line(blob2)

    add_text(
        slide, Inches(0), Inches(1.9), SLIDE_W, Inches(0.5),
        "FINAL THOUGHTS", size=13, bold=True, color=CYAN,
        align=PP_ALIGN.CENTER, font="Inter", letter_spacing=500,
    )

    add_gradient_text(
        slide, Inches(0), Inches(2.4), SLIDE_W, Inches(2.0),
        "Thank You",
        size=110, bold=True,
        stops=[(0, CYAN), (50, PURPLE_BRIGHT), (100, PINK)],
        angle_deg=15,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, Inches(0), Inches(4.9), SLIDE_W, Inches(0.6),
        "ShareBox 期望成為台灣社區鄰里最被信任的工具共享平台",
        size=22, bold=True, color=TEXT_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide, Inches(0), Inches(5.55), SLIDE_W, Inches(0.5),
        "讓閒置工具流動起來，減少浪費、連結鄰里",
        size=16, color=TEXT_SECONDARY,
        align=PP_ALIGN.CENTER,
    )

    # 分隔線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.17), Inches(6.3), Inches(1), Inches(0.04))
    _apply_gradient_fill(line, stops=[(0, CYAN, 100000), (100, PINK, 100000)], angle_deg=0)
    _no_line(line)

    add_text(
        slide, Inches(0), Inches(6.5), SLIDE_W, Inches(0.4),
        "SHAREBOX · 2026", size=12, bold=True, color=TEXT_MUTED,
        align=PP_ALIGN.CENTER, font="Inter", letter_spacing=600,
    )


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 10

    slide_cover(prs)
    slide_problem(prs, 2, total)
    slide_solution(prs, 3, total)
    slide_features(prs, 4, total)
    slide_market(prs, 5, total)
    slide_team(prs, 6, total)
    slide_timeline(prs, 7, total)
    slide_budget(prs, 8, total)
    slide_kpi(prs, 9, total)
    slide_closing(prs, 10, total)

    out = Path(__file__).parent / "專案簡報.pptx"
    prs.save(out)
    print(f"✓ saved → {out}  ({out.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    build()
