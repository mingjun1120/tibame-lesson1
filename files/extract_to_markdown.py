"""將 files/ 內各種格式文件抽取文字並輸出為結構化 Markdown 到 doc/。

支援格式：
- .docx  → python-docx，依段落樣式對應 Markdown 標題與清單
- .pptx  → python-pptx，每張投影片視為一節
- .pdf   → PyMuPDF，依字級推斷標題層級
- .md    → 直接複製
"""

from __future__ import annotations

import re
from collections import Counter
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent.parent
SRC_DIR = BASE_DIR / "files"
OUT_DIR = BASE_DIR / "doc"


# ---------- DOCX ----------

DOCX_STYLE_MAP = {
    "Title": "# ",
    "Heading 1": "## ",
    "Heading 2": "### ",
    "Heading 3": "#### ",
    "Heading 4": "##### ",
    "Heading 5": "###### ",
}


def _docx_iter_blocks(doc):
    """依文件順序產出 Paragraph / Table 物件。"""
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, doc)
        elif child.tag == qn("w:tbl"):
            yield Table(child, doc)


def _docx_table_to_md(table) -> list[str]:
    rows = [
        [cell.text.strip().replace("\n", " ") for cell in row.cells]
        for row in table.rows
    ]
    if not rows:
        return []
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = [
        "| " + " | ".join(rows[0]) + " |",
        "|" + "|".join(["---"] * width) + "|",
    ]
    for row in rows[1:]:
        out.append("| " + " | ".join(row) + " |")
    return out


def docx_to_markdown(path: Path) -> str:
    from docx.table import Table

    doc = Document(path)
    lines: list[str] = []
    prev_blank = True

    for block in _docx_iter_blocks(doc):
        if isinstance(block, Table):
            lines.append("")
            lines.extend(_docx_table_to_md(block))
            lines.append("")
            prev_blank = True
            continue

        para = block
        text = para.text.strip()
        style = para.style.name if para.style else "Normal"

        if not text:
            if not prev_blank:
                lines.append("")
                prev_blank = True
            continue

        if style in DOCX_STYLE_MAP:
            lines.append("")
            lines.append(f"{DOCX_STYLE_MAP[style]}{text}")
            lines.append("")
            prev_blank = True
        elif "List Bullet" in style or style.startswith("List Paragraph"):
            lines.append(f"- {text}")
            prev_blank = False
        elif "List Number" in style:
            lines.append(f"1. {text}")
            prev_blank = False
        else:
            lines.append(text)
            prev_blank = False

    return _normalize(lines)


# ---------- PPTX ----------

PPTX_CHROME_PATTERNS = [
    re.compile(r"^\d+\s*/\s*\d+$"),               # "02 / 10"
    re.compile(r"^SHAREBOX\s*·\s*專案簡報$", re.I),  # 頁尾
]

CJK_RE = re.compile(r"[一-鿿]")
EMOJI_ONLY_RE = re.compile(
    r"^[\U0001F300-\U0001FAFF☀-➿✀-➿⌀-⏿\s]+$"
)


def _is_chrome(text: str) -> bool:
    return any(pat.match(text) for pat in PPTX_CHROME_PATTERNS)


def _is_decorative(text: str) -> bool:
    """單字、純 emoji、純編號等視為裝飾性元素，排除於正文之外。"""
    if len(text) <= 1:
        return True
    if EMOJI_ONLY_RE.match(text):
        return True
    if re.fullmatch(r"[0-9]{1,2}", text):  # "01", "02"
        return True
    if re.fullmatch(r"[MK][0-9]+", text):   # "M1", "K2"
        return True
    return False


def _para_font_pt(para) -> float:
    """取段落第一個有設定字級的 run。"""
    for run in para.runs:
        if run.font.size:
            return run.font.size.pt
    return 0.0


def pptx_to_markdown(path: Path) -> str:
    prs = Presentation(path)
    lines: list[str] = ["# " + path.stem, ""]

    for idx, slide in enumerate(prs.slides, start=1):
        items: list[tuple[float, float, str, float]] = []  # (top, left, text, size)
        tables_md: list[list[str]] = []

        for shape in slide.shapes:
            try:
                top = shape.top or 0
                left = shape.left or 0
            except Exception:
                top, left = 0, 0

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text or _is_chrome(text):
                        continue
                    size = _para_font_pt(para)
                    items.append((top, left, text, size))
            elif shape.has_table:
                table = shape.table
                rows = [
                    [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    for row in table.rows
                ]
                if rows:
                    width = max(len(r) for r in rows)
                    rows = [r + [""] * (width - len(r)) for r in rows]
                    md = [
                        "| " + " | ".join(rows[0]) + " |",
                        "|" + "|".join(["---"] * width) + "|",
                    ]
                    for row in rows[1:]:
                        md.append("| " + " | ".join(row) + " |")
                    tables_md.append(md)

        # 過濾裝飾性元素
        items = [it for it in items if not _is_decorative(it[2])]

        if not items and not tables_md:
            lines.append(f"## 投影片 {idx}")
            lines.append("")
            continue

        # 主標題挑選策略：上半部 + 含中文 + 字級較大者優先
        slide_height = (prs.slide_height or 1)
        upper_half = [it for it in items if it[0] < slide_height * 0.45]
        title_pool = [it for it in upper_half if CJK_RE.search(it[2]) and 3 <= len(it[2]) <= 30]
        if not title_pool:
            title_pool = [it for it in items if CJK_RE.search(it[2]) and 3 <= len(it[2]) <= 30]

        title_item = None
        if title_pool:
            # 字級大者優先；同字級時取位置靠上者
            title_item = max(title_pool, key=lambda x: (x[3], -x[0]))

        title = title_item[2] if title_item else f"投影片 {idx}"
        lines.append(f"## 投影片 {idx}：{title}")
        lines.append("")

        # 其餘文字依視覺位置（上→下、左→右）列出
        body_items = [it for it in items if it is not title_item]
        body_items.sort(key=lambda x: (round(x[0] / 100000), x[1]))

        for top, left, text, size in body_items:
            lines.append(f"- {text}")
        lines.append("")

        for md in tables_md:
            lines.extend(md)
            lines.append("")

    return _normalize(lines)


# ---------- PDF ----------

HEADING_MAX_LEN = 60  # 超過這個字數即使字級大也不視為標題


def _bbox_overlap(a, b) -> bool:
    """檢查兩個 bbox (x0,y0,x1,y1) 是否有顯著重疊。"""
    ax0, ay0, ax1, ay1 = a
    bx0, by0, bx1, by1 = b
    # 只要 y 軸有交集且 x 軸有交集即視為在表格範圍內
    return not (ax1 < bx0 or bx1 < ax0 or ay1 < by0 or by1 < ay0)


def _table_to_markdown(extracted_rows) -> list[str]:
    """將 PyMuPDF Table.extract() 結果轉成 Markdown 表格。"""
    rows = [
        [(cell or "").strip().replace("\n", " ") for cell in row]
        for row in extracted_rows
        if any((cell or "").strip() for cell in row)
    ]
    if not rows:
        return []
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    out = [
        "| " + " | ".join(rows[0]) + " |",
        "|" + "|".join(["---"] * width) + "|",
    ]
    for row in rows[1:]:
        out.append("| " + " | ".join(row) + " |")
    return out


def _dominant_body_size(doc) -> float:
    """以字數加權，找出最常見的字級當作本文字級。"""
    weights: Counter = Counter()
    for page in doc:
        for block in page.get_text("dict")["blocks"]:
            if block.get("type") != 0:
                continue
            for line in block["lines"]:
                for span in line["spans"]:
                    txt = span["text"].strip()
                    if txt:
                        weights[round(span["size"], 1)] += len(txt)
    if not weights:
        return 10.0
    return weights.most_common(1)[0][0]


def pdf_to_markdown(path: Path) -> str:
    doc = fitz.open(path)
    body_size = _dominant_body_size(doc)
    h1_threshold = body_size * 1.8
    h2_threshold = body_size * 1.3
    h3_threshold = body_size * 1.1

    def classify(size: float, text: str) -> str:
        if len(text) > HEADING_MAX_LEN:
            return ""
        if size >= h1_threshold:
            return "# "
        if size >= h2_threshold:
            return "## "
        if size >= h3_threshold and len(text) < 40:
            return "### "
        return ""

    lines: list[str] = []

    for page in doc:
        # 表格偵測：先記下表格 bbox，產出對應 Markdown，並排除文字輸出
        try:
            tables = page.find_tables()
        except Exception:
            tables = None

        table_blocks = []  # [(bbox, markdown_lines)]
        if tables:
            for tbl in tables.tables:
                try:
                    md = _table_to_markdown(tbl.extract())
                except Exception:
                    md = []
                if md:
                    table_blocks.append((tbl.bbox, md))

        rendered_tables = [False] * len(table_blocks)

        # 將每個文字 block 抽成 (y0, level_prefix, text)
        items: list[tuple[float, str, str, tuple]] = []

        for block in page.get_text("dict")["blocks"]:
            if block.get("type") != 0:
                continue
            bb = block["bbox"]

            # 若 block 落在表格範圍內，跳過（表格會另外輸出）
            in_table = False
            for ti, (tbbox, _) in enumerate(table_blocks):
                if _bbox_overlap(bb, tbbox):
                    in_table = True
                    break
            if in_table:
                continue

            # 同 block 內合併所有 line（同段文字），取最大字級
            block_text_parts: list[str] = []
            block_max_size = 0.0
            for line in block["lines"]:
                line_text = "".join(span["text"] for span in line["spans"]).strip()
                if not line_text:
                    continue
                block_text_parts.append(line_text)
                for span in line["spans"]:
                    if span["text"].strip():
                        block_max_size = max(block_max_size, span["size"])

            if not block_text_parts:
                continue

            # 中文段落：行與行間直接相連（避免破詞）；其它語言以空白連接
            text = "".join(block_text_parts)
            prefix = classify(block_max_size, text)
            items.append((bb[1], prefix, text, bb))

        # 把表格也加入 items，依 y0 排序
        for ti, (tbbox, md) in enumerate(table_blocks):
            items.append((tbbox[1], "__TABLE__", "\n".join(md), tbbox))

        items.sort(key=lambda x: x[0])

        last_kind = None
        for y0, prefix, text, bb in items:
            if prefix == "__TABLE__":
                lines.append("")
                lines.extend(text.split("\n"))
                lines.append("")
                last_kind = "table"
                continue

            if prefix:
                lines.append("")
                lines.append(f"{prefix}{text}")
                lines.append("")
                last_kind = "heading"
                continue

            # bullet 偵測（中文常見符號）
            bullet_match = re.match(r"^\s*[•●◆▪■□・·‧]\s*(.+)", text)
            num_match = re.match(r"^\s*(\d+[\.、)])\s*(.+)", text)

            if bullet_match:
                lines.append(f"- {bullet_match.group(1).strip()}")
                last_kind = "list"
            elif num_match and len(text) < 80:
                lines.append(f"{num_match.group(1)} {num_match.group(2).strip()}")
                last_kind = "list"
            else:
                if last_kind in {"list", "table"}:
                    lines.append("")
                # 處理段落內以「・」開頭的多個項目（一個 block 多個 bullet）
                if "・" in text and text.count("・") >= 2:
                    parts = [p.strip() for p in text.split("・") if p.strip()]
                    for p in parts:
                        lines.append(f"- {p}")
                    last_kind = "list"
                else:
                    lines.append(text)
                    last_kind = "text"

        lines.append("")

    return _normalize(lines)


# ---------- Markdown 既有檔 ----------

def md_passthrough(path: Path) -> str:
    return path.read_text(encoding="utf-8")


# ---------- 工具 ----------

def _normalize(lines: list[str]) -> str:
    """合併連續空行、去除尾端空白。"""
    out: list[str] = []
    blank = False
    for ln in lines:
        ln = ln.rstrip()
        if ln == "":
            if not blank:
                out.append("")
            blank = True
        else:
            out.append(ln)
            blank = False
    # 去除開頭多餘空行
    while out and out[0] == "":
        out.pop(0)
    if not out or out[-1] != "":
        out.append("")
    return "\n".join(out)


EXTRACTORS = {
    ".docx": docx_to_markdown,
    ".pptx": pptx_to_markdown,
    ".pdf": pdf_to_markdown,
    ".md": md_passthrough,
}


def main() -> None:
    OUT_DIR.mkdir(exist_ok=True)
    processed = []

    for src in sorted(SRC_DIR.iterdir()):
        if src.is_dir() or src.name.startswith("."):
            continue
        if src.suffix.lower() not in EXTRACTORS:
            continue

        extractor = EXTRACTORS[src.suffix.lower()]
        try:
            content = extractor(src)
        except Exception as exc:  # noqa: BLE001
            print(f"[警告] 處理 {src.name} 失敗：{exc}")
            continue

        out_path = OUT_DIR / f"{src.stem}.md"
        out_path.write_text(content, encoding="utf-8")
        processed.append((src.name, out_path.name, len(content)))
        print(f"✓ {src.name:<40} → doc/{out_path.name}  ({len(content)} chars)")

    print(f"\n完成：共輸出 {len(processed)} 份檔案至 {OUT_DIR}")


if __name__ == "__main__":
    main()
